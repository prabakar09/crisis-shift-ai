import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Palette
    C_BG = RGBColor(10, 14, 23)           # #0A0E17 Deep Obsidian
    C_CARD = RGBColor(19, 27, 44)         # #131B2C Slate Navy
    C_BORDER = RGBColor(39, 52, 79)       # #27344F Subtle Border
    C_AMBER = RGBColor(245, 158, 11)      # #F59E0B Cinema Gold
    C_CYAN = RGBColor(56, 189, 248)       # #38BDF8 Satellite Cyan
    C_WHITE = RGBColor(248, 250, 252)     # #F8FAFC Bright White
    C_MUTED = RGBColor(148, 163, 184)     # #94A3B8 Silver Gray
    C_GREEN = RGBColor(16, 185, 129)      # #10B981 Emerald Green
    C_RED = RGBColor(239, 68, 68)         # #EF4444 Alert Crimson

    def add_slide():
        s = prs.slides.add_slide(blank_layout)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.fill.background()
        return s

    def add_header(s, category, title, subtitle=None):
        tb_cat = s.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = tb_cat.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = C_AMBER

        tb_title = s.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.65))
        tf_t = tb_title.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = C_WHITE

        if subtitle:
            tb_sub = s.shapes.add_textbox(Inches(0.8), Inches(1.42), Inches(11.7), Inches(0.4))
            tf_s = tb_sub.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(13)
            p_s.font.color.rgb = C_MUTED

    def make_card(s, left, top, width, height, title, items, badge=None, accent=C_CYAN):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        c.fill.solid()
        c.fill.fore_color.rgb = C_CARD
        c.line.color.rgb = C_BORDER
        c.line.width = Pt(1.2)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.28)
        tf.margin_right = Inches(0.28)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.2)

        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = f"[{badge}]"
            p_badge.font.size = Pt(10)
            p_badge.font.bold = True
            p_badge.font.color.rgb = accent
            p_badge.space_after = Pt(4)
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]

        p_title.text = title
        p_title.font.size = Pt(17)
        p_title.font.bold = True
        p_title.font.color.rgb = C_WHITE
        p_title.space_after = Pt(10)

        for item in items:
            p = tf.add_paragraph()
            p.space_before = Pt(3)
            p.space_after = Pt(3)
            if ":" in item:
                lead, trail = item.split(":", 1)
                r0 = p.add_run()
                r0.text = "• " + lead + ":"
                r0.font.bold = True
                r0.font.color.rgb = C_WHITE
                r0.font.size = Pt(12)

                r1 = p.add_run()
                r1.text = trail
                r1.font.bold = False
                r1.font.color.rgb = C_MUTED
                r1.font.size = Pt(12)
            else:
                r0 = p.add_run()
                r0.text = "• " + item
                r0.font.color.rgb = C_MUTED
                r0.font.size = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = add_slide()

    tag = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(5.2), Inches(0.42))
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(26, 36, 56)
    tag.line.color.rgb = C_AMBER
    tag.line.width = Pt(1)
    tag_tf = tag.text_frame
    tag_p = tag_tf.paragraphs[0]
    tag_p.text = "GOOGLE GEMINI 3.8 ENTERPRISE AGENT PLATFORM"
    tag_p.font.size = Pt(10)
    tag_p.font.bold = True
    tag_p.font.color.rgb = C_AMBER
    tag_p.alignment = PP_ALIGN.CENTER

    tb_main = s1.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.3))
    tf1 = tb_main.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "CrisisShift OS"
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Autonomous Cinema Command & Set Orchestration"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = C_CYAN
    p2.space_before = Pt(6)

    tb_desc = s1.shapes.add_textbox(Inches(0.8), Inches(4.35), Inches(11.7), Inches(1.1))
    tf_d = tb_desc.text_frame
    tf_d.word_wrap = True
    pd = tf_d.paragraphs[0]
    pd.text = "Transforming real-world film set chaos into deterministic autonomous cinema decisions with Gemini 3.8 Multi-Agent Swarm, Real-Time Satellite Atmospheric Radar, and Surface Asphalt Friction Physics."
    pd.font.size = Pt(15)
    pd.font.color.rgb = C_MUTED

    meta_items = [
        ("TRACK", "Agentic Cinema: Director Track", C_AMBER),
        ("AI ENGINE", "Gemini 3.8 Flash Swarm", C_CYAN),
        ("LOCALIZATION", "Bilingual Tanglish / Regional", C_GREEN),
        ("ESTIMATED ROI", "$50k - $500k Saved / Rainout", C_WHITE)
    ]
    for i, (m_tag, m_val, m_c) in enumerate(meta_items):
        box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 2.95), Inches(5.8), Inches(2.8), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = C_CARD
        box.line.color.rgb = C_BORDER
        btf = box.text_frame
        btf.margin_top = Inches(0.12)
        bp0 = btf.paragraphs[0]
        bp0.text = m_tag
        bp0.font.size = Pt(9)
        bp0.font.bold = True
        bp0.font.color.rgb = m_c
        bp1 = btf.add_paragraph()
        bp1.text = m_val
        bp1.font.size = Pt(12)
        bp1.font.bold = True
        bp1.font.color.rgb = C_WHITE

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary
    # -------------------------------------------------------------
    s2 = add_slide()
    add_header(s2, "Executive Overview", "CrisisShift OS at a Glance", 
               "Bridging satellite physics and production logistics into an autonomous command system.")

    make_card(s2, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "The Core Problem",
              [
                  "Astronomical Cost: Outdoor film shoots lose $50,000 to $500,000 (Rs 40 Lakhs - 4 Crores) per cancelled day.",
                  "Life-Threatening Hazards: Stunt cars hydroplane on wet tarmac; 5-ton Technocranes tip over in rain-softened soil.",
                  "Static Disconnection: Paper call sheets and generic phone apps fail to provide mathematical risk arbitration."
              ], "PAIN POINT", C_RED)

    make_card(s2, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "The Innovation",
              [
                  "Deterministic Swarm: Gemini 3.8 arbitrates telemetry across 4 specialized autonomous agents.",
                  "Physical Grounding: Real-time radar plus dynamic friction (mu) and soil bearing capacity (sigma).",
                  "Cultural Precision: First-of-its-kind bilingual Tanglish dispatch engine for regional film crew communication."
              ], "SOLUTION", C_CYAN)

    make_card(s2, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "The Business Impact",
              [
                  "Instant Execution: Sub-15s autonomous verdict replacing hours of costly set debates.",
                  "Completion Bond Discounts: 15-25% lower insurance premiums through verified physics audit trails.",
                  "High Margin SaaS: 85%+ gross profit margin across Indie, Studio, and Enterprise production tiers."
              ], "VALUE CREATION", C_GREEN)

    # -------------------------------------------------------------
    # SLIDE 3: Problem Statement
    # -------------------------------------------------------------
    s3 = add_slide()
    add_header(s3, "01. Problem Statement", "Multi-Million Dollar Losses & Safety Blindspots",
               "Major productions worldwide face catastrophic financial bleeds and life safety hazards on unverified terrain.")

    make_card(s3, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "Financial Catastrophe",
              [
                  "Daily Burn Rate: Film productions bleed $50,000 to $500,000 every single day an outdoor shoot is rained out.",
                  "Idle Sunk Costs: Star cast date locks, 200+ union crew standby wages, and Technocrane daily lease rentals run regardless.",
                  "Budget Blowouts: Just 3 unpredicted weather cancellations can swell overall project budget by over 20-30%."
              ], "FINANCIAL BLEED", C_RED)

    make_card(s3, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "Fatal Stunt & Rig Hazards",
              [
                  "Hydroplaning Wrecks: High-speed chase stunts on damp asphalt suffer fatal skids when friction coefficient drops below 0.35.",
                  "Technocrane Overturns: 50-foot camera cranes exert extreme outrigger pressure, risking collapse into rain-softened mud.",
                  "Zero Sensory Data: Directors and 1st ADs cannot visually estimate road friction or subsurface soil density."
              ], "LIFE SAFETY", C_AMBER)

    make_card(s3, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "Failure of Existing Tools",
              [
                  "Static Call Sheets: Paper call sheets and Excel matrices are outdated the moment micro-weather shifts.",
                  "Consumer Weather Apps: Generic phone apps lack precision atmospheric radar or camera rig engineering context.",
                  "Language Bottlenecks: Regional film crews (Tamil, Hindi, Telugu) struggle with delayed formal English directives."
              ], "OPERATIONAL BREAKDOWN", C_CYAN)

    # -------------------------------------------------------------
    # SLIDE 4: The Solution & What We Built
    # -------------------------------------------------------------
    s4 = add_slide()
    add_header(s4, "02. The Solution", "CrisisShift OS Autonomous Command Architecture",
               "An end-to-end intelligent cockpit replacing guesswork with autonomous multi-agent orchestration.")

    # 4 Autonomous Agents
    make_card(s4, Inches(0.8), Inches(2.0), Inches(2.75), Inches(4.8),
              "Agent 1: Orbital Scout",
              [
                  "Radar Telemetry: Live Open-Meteo atmospheric radar.",
                  "Precipitation & Wind: Instant wind vector and rain volume extraction.",
                  "ArcGIS Recon: High-resolution satellite optical imagery.",
                  "3D Pydeck Radar: Live location beacon."
              ], "SATELLITE", C_CYAN)

    make_card(s4, Inches(3.78), Inches(2.0), Inches(2.75), Inches(4.8),
              "Agent 2: Terrain Physics",
              [
                  "Asphalt Friction (mu): Computes road grip (0.32 wet vs 0.85 dry).",
                  "Soil Bearing (sigma): Checks crane outrigger load pressure against structural limits.",
                  "Stunt Safety Index: Dynamic braking distance multipliers."
              ], "PHYSICS", C_AMBER)

    make_card(s4, Inches(6.76), Inches(2.0), Inches(2.75), Inches(4.8),
              "Agent 3: Gemini 3.8 Swarm",
              [
                  "Deterministic Verdict: Strict mathematical Go/No-Go decision.",
                  "RED HALT Enforcement: If Rain > 1.0mm or mu < 0.40, immediate weather halt.",
                  "Indoor Relocation: Automatic stage contingency routing.",
                  "Zero Hallucination: Grounded in sensor logic."
              ], "ORCHESTRATOR", C_GREEN)

    make_card(s4, Inches(9.75), Inches(2.0), Inches(2.75), Inches(4.8),
              "Agent 4: Crew Dispatcher",
              [
                  "Tanglish First: Native Tamil+English WhatsApp alerts.",
                  "Multilingual Support: Hindi, Telugu, and English dispatches.",
                  "1-Click PDF Export: FPDF2 Unicode audit dispatches.",
                  "SQLite Audit Trail: Permanent database logging."
              ], "DISPATCH", C_WHITE)

    # -------------------------------------------------------------
    # SLIDE 5: How It Works (Workflow Pipeline)
    # -------------------------------------------------------------
    s5 = add_slide()
    add_header(s5, "03. How It Works", "5-Stage Deterministic Production Pipeline",
               "From location coordinates to emergency crew WhatsApp broadcast in under 15 seconds.")

    steps = [
        ("01", "Location & Scene Metadata Ingestion", 
         "Director or 1st AD inputs shoot GPS coordinates, scheduled call time, and scene parameters (e.g. Highway Car Chase, Technocrane Rig).", C_CYAN),
        ("02", "Orbital Satellite Radar Recon", 
         "Agent 1 queries OpenStreetMap Nominatim and Open-Meteo atmospheric radar, pulling real-time rain volume, wind gusts, and cloud cover.", C_AMBER),
        ("03", "Terrain Physics Computation", 
         "Agent 2 calculates surface asphalt friction (mu = 0.32 - 0.85) and tests Technocrane outrigger ground pressure against soil bearing limits.", C_GREEN),
        ("04", "Gemini 3.8 Swarm Decision Arbitration", 
         "Agent 3 synthesizes telemetry. If Rain > 1.0mm or mu < 0.40, it triggers an immediate RED HALT and routes crew to covered Soundstage B.", C_RED),
        ("05", "Bilingual Tanglish Crew Dispatch & Audit Log", 
         "Agent 4 formats WhatsApp alerts in native Tanglish/English for instant broadcast, auto-logging the full compliance record to SQLite and PDF.", C_WHITE)
    ]

    for i, (num, stitle, sdesc, scol) in enumerate(steps):
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.95 + i * 1.02), Inches(11.7), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = C_CARD
        box.line.color.rgb = C_BORDER
        box.line.width = Pt(1)

        # Number badge
        nb = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(2.1 + i * 1.02), Inches(0.9), Inches(0.6))
        nb.fill.solid()
        nb.fill.fore_color.rgb = RGBColor(26, 36, 56)
        nb.line.color.rgb = scol
        ntf = nb.text_frame
        np = ntf.paragraphs[0]
        np.text = num
        np.font.size = Pt(13)
        np.font.bold = True
        np.font.color.rgb = scol
        np.alignment = PP_ALIGN.CENTER

        # Content text
        tb = s5.shapes.add_textbox(Inches(2.05), Inches(2.02 + i * 1.02), Inches(10.2), Inches(0.78))
        ttf = tb.text_frame
        ttf.word_wrap = True
        ttf.margin_left = ttf.margin_top = ttf.margin_right = ttf.margin_bottom = 0
        tp1 = ttf.paragraphs[0]
        tp1.text = stitle
        tp1.font.size = Pt(13)
        tp1.font.bold = True
        tp1.font.color.rgb = C_WHITE
        tp2 = ttf.add_paragraph()
        tp2.text = sdesc
        tp2.font.size = Pt(11)
        tp2.font.color.rgb = C_MUTED
        tp2.space_before = Pt(2)

    # -------------------------------------------------------------
    # SLIDE 6: Core Benefits
    # -------------------------------------------------------------
    s6 = add_slide()
    add_header(s6, "04. Core Benefits", "Why Film Productions Adopt CrisisShift OS",
               "Speed, life safety assurance, and flawless set coordination across departments.")

    make_card(s6, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "Sub-15s Latency",
              [
                  "Lightning Fast: Orbital radar, friction math, and swarm decision execute in under 15 seconds.",
                  "Replaces Endless Debates: Eliminates frantic phone calls between director, producer, and camera department in the rain.",
                  "Real-Time Adaptability: Automatically monitors micro-climate shifts and updates production dispatches."
              ], "OPERATIONAL SPEED", C_CYAN)

    make_card(s6, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "Zero Stunt Casualties",
              [
                  "Mathematical Rigor: Eliminates subjective guesses by calculating true braking distance and grip coefficients.",
                  "Rig & Gear Protection: Prevents $500,000 Technocranes and luxury cinema camera packages from mud sinking or tip-overs.",
                  "Insurance Compliance: Generates timestamped PDF dispatches for completion bond underwriting and liability audits."
              ], "LIFE SAFETY", C_GREEN)

    make_card(s6, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.8),
              "Native Tanglish Edge",
              [
                  "Cultural Precision: Local crew members in Chennai, Hyderabad, and Mumbai understand Tanglish immediately.",
                  "Actionable On-Set Language: Direct, colloquial directives (e.g. 'Vandi slow pannunga, technocrane plates podunga').",
                  "WhatsApp Broadcast: Formatted for 1-click copy-paste into departmental chat groups without delay."
              ], "LOCALIZATION", C_AMBER)

    # -------------------------------------------------------------
    # SLIDE 7: Business Value, ROI & Profit Model
    # -------------------------------------------------------------
    s7 = add_slide()
    add_header(s7, "05. Business Value & Profit", "Financial Impact, ROI & High-Margin Revenue Model",
               "Saving massive production budgets while unlocking an enterprise SaaS revenue stream.")

    make_card(s7, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8),
              "Direct Financial ROI for Productions",
              [
                  "Rainout Prevention: Saves $50,000 to $500,000 (Rs 40 Lakhs - Rs 4 Crores) per averted cancellation day by shifting crew indoors early.",
                  "15-25% Insurance Discounts: Completion bond insurers provide lower policy premiums for productions utilizing verified physics risk monitoring.",
                  "Gear Wreckage Avoidance: Prevents $200k+ cinema camera bodies and specialized stunt vehicles from crash damage.",
                  "Crew Overtime Optimization: Eliminates paid standby hours for 200+ crew by orchestrating immediate soundstage fallback."
              ], "PRODUCTION ROI", C_GREEN)

    make_card(s7, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8),
              "Commercial Monetization & Profit Model",
              [
                  "Indie Tier (Rs 25,000 / month): Targeted at ad commercials and indie filmmakers for basic atmospheric radar and physics.",
                  "Blockbuster Tier (Rs 1,50,000 / shoot): Full 4-agent swarm, Tanglish WhatsApp dispatcher, and PDF audit engine for 90-day schedules.",
                  "Enterprise Studio Slate ($25,000 / year): Unlimited access for major production banners (Lyca, Sun Pictures, Dharma, Netflix India).",
                  "85%+ Gross Profit Margin: High-efficiency Gemini 3.8 Flash API inference provides near-zero marginal cost per recon run."
              ], "COMMERCIAL MODEL", C_AMBER)

    # -------------------------------------------------------------
    # SLIDE 8: Competitive Advantage
    # -------------------------------------------------------------
    s8 = add_slide()
    add_header(s8, "06. Market Comparison", "Why CrisisShift OS Dominates Traditional Tools",
               "Comparing traditional production tools against CrisisShift OS Autonomous Cinema Command.")

    col_w = [Inches(3.2), Inches(2.8), Inches(2.8), Inches(2.9)]
    row_h = Inches(0.82)
    left_start = Inches(0.8)
    top_start = Inches(2.0)

    headers = ["Feature / Capability", "Static Call Sheets", "Weather Apps", "CrisisShift OS"]
    rows = [
        ["Real-Time Radar & Satellite", "None / Paper Only", "City-Level Forecast", "Live Open-Meteo + ArcGIS Recon"],
        ["Surface Friction (mu) Physics", "None", "None", "Exact Mathematical Model"],
        ["Autonomous Decision Making", "Manual Set Debates", "User Must Interpret", "Deterministic Gemini 3.8 Swarm"],
        ["Set Localization (Tanglish)", "English / Manual", "English Only", "Native Tanglish/Regional Broadcast"],
        ["1-Click PDF Audit Records", "Manual Paperwork", "None", "Instant FPDF2 Unicode PDF Export"]
    ]

    for c_idx, h_text in enumerate(headers):
        c_x = left_start + sum(col_w[:c_idx])
        cell = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_x, top_start, col_w[c_idx], Inches(0.6))
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(26, 36, 56) if c_idx < 3 else C_AMBER
        cell.line.color.rgb = C_BORDER
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_BG if c_idx == 3 else C_WHITE
        p.alignment = PP_ALIGN.CENTER

    for r_idx, row in enumerate(rows):
        r_y = top_start + Inches(0.65) + r_idx * row_h
        for c_idx, val in enumerate(row):
            c_x = left_start + sum(col_w[:c_idx])
            cell = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_x, r_y, col_w[c_idx], row_h - Inches(0.08))
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD if c_idx < 3 else RGBColor(22, 34, 52)
            cell.line.color.rgb = C_BORDER if c_idx < 3 else C_CYAN
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = C_WHITE if c_idx == 0 or c_idx == 3 else C_MUTED
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 9: Vision & Roadmap
    # -------------------------------------------------------------
    s9 = add_slide()
    add_header(s9, "07. Roadmap & Vision", "Scaling to the Global Cinema Operating System",
               "Evolving from single set orchestrations to studio-wide enterprise operations.")

    make_card(s9, Inches(0.8), Inches(2.0), Inches(3.75), Inches(4.8),
              "Hackathon Launch (Current)",
              [
                  "Gemini 3.8 Swarm: 4-agent autonomous directorial logic.",
                  "Radar & Physics Engine: Live Open-Meteo, asphalt friction, and crane load math.",
                  "Bilingual Localization: Tanglish WhatsApp dispatcher.",
                  "Audit Compliance: SQLite database and FPDF2 PDF exports."
              ], "PHASE 1: CURRENT", C_GREEN)

    make_card(s9, Inches(4.78), Inches(2.0), Inches(3.75), Inches(4.8),
              "Studio MCP Integration",
              [
                  "MCP Server Adapter: Model Context Protocol sync with Movie Magic Scheduling.",
                  "Live Drone Video Feeds: Direct aerial optical reconnaissance.",
                  "Twilio WhatsApp API: Automated simultaneous dispatch to 200+ crew phones.",
                  "Multi-Camera Tracking: Dynamic lighting and cloud shadow calculation."
              ], "PHASE 2: NEAR-TERM", C_CYAN)

    make_card(s9, Inches(8.75), Inches(2.0), Inches(3.75), Inches(4.8),
              "Global Cinema OS Standard",
              [
                  "Multi-Unit Orchestration: Autonomous sync between 1st Unit and 2nd Unit stunt teams.",
                  "AI Stunt Simulations: Pre-visualized trajectory and friction collision modeling.",
                  "Automated Insurance Filing: Instant claim settlement dispatch for weather delays.",
                  "Enterprise Studio Cockpit: Global dashboard for multi-film slates."
              ], "PHASE 3: LONG-TERM", C_AMBER)

    # -------------------------------------------------------------
    # SLIDE 10: Conclusion & Call to Action
    # -------------------------------------------------------------
    s10 = add_slide()

    tag10 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(3.8), Inches(0.42))
    tag10.fill.solid()
    tag10.fill.fore_color.rgb = RGBColor(26, 36, 56)
    tag10.line.color.rgb = C_AMBER
    tag10.line.width = Pt(1)
    t10_tf = tag10.text_frame
    t10_p = t10_tf.paragraphs[0]
    t10_p.text = "THE BLOCKBUSTER HACKATHON"
    t10_p.font.size = Pt(10)
    t10_p.font.bold = True
    t10_p.font.color.rgb = C_AMBER
    t10_p.alignment = PP_ALIGN.CENTER

    tb_end = s10.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.2))
    tf_e = tb_end.text_frame
    tf_e.word_wrap = True
    pe1 = tf_e.paragraphs[0]
    pe1.text = "From Chaos on Set to Autonomous Mastery."
    pe1.font.size = Pt(44)
    pe1.font.bold = True
    pe1.font.color.rgb = C_WHITE

    pe2 = tf_e.add_paragraph()
    pe2.text = "CrisisShift OS protects crew lives, saves multi-crore film budgets, and brings deterministic AI command to global cinema."
    pe2.font.size = Pt(18)
    pe2.font.bold = True
    pe2.font.color.rgb = C_CYAN
    pe2.space_before = Pt(8)

    # 3 Summary Pillar Cards
    summary_pillars = [
        ("Sub-15s Speed", "Instant satellite recon and physics-backed decisions.", C_CYAN),
        ("Rs 40L - Rs 4Cr Saved", "Proven cost avoidance per rainout day.", C_GREEN),
        ("Gemini 3.8 Swarm", "The ultimate autonomous director for modern cinema.", C_AMBER)
    ]
    for i, (p_title, p_desc, p_col) in enumerate(summary_pillars):
        p_box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i * 3.98), Inches(4.5), Inches(3.75), Inches(1.8))
        p_box.fill.solid()
        p_box.fill.fore_color.rgb = C_CARD
        p_box.line.color.rgb = p_col
        p_box.line.width = Pt(1.2)
        ptf = p_box.text_frame
        ptf.margin_top = Inches(0.2)
        ptf.margin_left = Inches(0.25)
        ptf.margin_right = Inches(0.25)
        pp0 = ptf.paragraphs[0]
        pp0.text = p_title
        pp0.font.size = Pt(18)
        pp0.font.bold = True
        pp0.font.color.rgb = p_col
        pp1 = ptf.add_paragraph()
        pp1.text = p_desc
        pp1.font.size = Pt(13)
        pp1.font.color.rgb = C_WHITE
        pp1.space_before = Pt(6)

    output_path = "CrisisShift_OS_Presentation.pptx"
    prs.save(output_path)
    print(f"Deck saved to: {output_path}")

if __name__ == "__main__":
    build_deck()
